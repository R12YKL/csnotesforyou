import os
import json
from pptx import Presentation

notes_dir = 'public/Notes'
output_file = 'public/data/notes.json'

# Mapping of file numbers to section categories
section_mapping = {
    1: "Fundamentals", 2: "Fundamentals", 3: "Fundamentals", 26: "Fundamentals",
    4: "Hardware & Systems", 5: "Hardware & Systems", 16: "Hardware & Systems", 17: "Hardware & Systems",
    6: "Python Programming", 7: "Python Programming", 8: "Python Programming", 9: "Python Programming",
    10: "Python Programming", 11: "Python Programming", 13: "Python Programming", 14: "Python Programming",
    15: "Python Programming", 19: "Python Programming", 20: "Python Programming",
    12: "Problem Solving & Algorithms", 23: "Problem Solving & Algorithms",
    21: "Web & Databases", 22: "Web & Databases", 27: "Web & Databases",
    18: "Advanced Topics", 24: "Advanced Topics", 25: "Advanced Topics"
}

# Extract content from PPTX files
notes_data = {}
files_processed = 0

for filename in sorted(os.listdir(notes_dir)):
    if filename.endswith('.pptx'):
        file_path = os.path.join(notes_dir, filename)
        
        try:
            # Extract file number from filename
            file_num = int(filename.split(' - ')[0].split('  ')[0].strip())
            
            # Get section from mapping
            section = section_mapping.get(file_num, "Advanced Topics")
            
            # Initialize section if needed
            if section not in notes_data:
                notes_data[section] = []
            
            # Extract presentation
            prs = Presentation(file_path)
            
            # Get title from filename
            title = filename.replace('.pptx', '').replace(f'{file_num} - ', '').replace(f'{file_num}  ', '').strip()
            
            # Extract all slide content
            slides_content = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides_content.append({
                        "slide": slide_idx,
                        "content": " ".join(slide_text)
                    })
            
            # Combine all slide content for display
            full_content = "\n\n".join([f"**Slide {s['slide']}:** {s['content']}" for s in slides_content])
            
            # Add to notes
            notes_data[section].append({
                "id": str(file_num),
                "title": title,
                "content": full_content if full_content else f"Presentation: {title}",
                "file": f"Notes/{filename}"
            })
            
            files_processed += 1
            print(f"✓ Processed: {filename}")
            
        except Exception as e:
            print(f"✗ Error processing {filename}: {str(e)}")

# Convert to expected format with sections array
final_notes = {
    "sections": [
        {
            "title": section_name,
            "notes": notes_data.get(section_name, [])
        }
        for section_name in [
            "Fundamentals",
            "Hardware & Systems", 
            "Python Programming",
            "Problem Solving & Algorithms",
            "Web & Databases",
            "Advanced Topics"
        ]
    ]
}

# Save to JSON
with open(output_file, 'w', encoding='utf-8') as f:
    json.dump(final_notes, f, indent=2, ensure_ascii=False)

print(f"\n✓ Successfully processed {files_processed} presentations")
print(f"✓ Saved to {output_file}")
